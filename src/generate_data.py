import argparse
import os
import random
import uuid
from dataclasses import dataclass
from datetime import datetime, timedelta, timezone

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from dateutil.relativedelta import relativedelta
from faker import Faker
from pptx import Presentation
from pptx.util import Inches, Pt


UTC = timezone.utc


def month_window(month_yyyy_mm: str) -> tuple[datetime, datetime]:
    start = datetime.fromisoformat(month_yyyy_mm + "-01").replace(tzinfo=UTC)
    end = (start + relativedelta(months=1))
    return start, end


def ensure_dir(path: str) -> None:
    os.makedirs(path, exist_ok=True)


def weighted_choice(rng: random.Random, items: list[str], weights: list[float]) -> str:
    return rng.choices(items, weights=weights, k=1)[0]


def clamp_dt(dt: datetime, start: datetime, end: datetime) -> datetime:
    if dt < start:
        return start
    if dt >= end:
        return end - timedelta(seconds=1)
    return dt


def ts_iso(dt: datetime) -> str:
    return dt.astimezone(UTC).isoformat().replace("+00:00", "Z")


@dataclass(frozen=True)
class Config:
    n_users: int = 1000
    n_txns: int = 5000
    n_events: int = 10000
    # anomaly rates (purposely different from the examples in the brief)
    rate_txn_invalid_amount: float = 0.01  # <=0 or too large
    rate_txn_orphan_receiver: float = 0.01  # receiver_user_id not found in users
    rate_txn_outside_month: float = 0.02  # timestamp outside reporting month (but still after signup)
    rate_event_unknown_user: float = 0.02  # user_id not in users


def generate_users(rng: random.Random, fake: Faker, month_start: datetime, month_end: datetime, n_users: int) -> pd.DataFrame:
    countries = ["GB", "US", "ES", "BR", "NG", "DE", "FR", "IN", "MX", "PH"]
    country_w = [0.18, 0.18, 0.08, 0.08, 0.08, 0.1, 0.08, 0.1, 0.06, 0.06]
    oses = ["ios", "android", "web"]
    os_w = [0.45, 0.45, 0.10]
    kyc = ["none", "pending", "approved", "rejected"]
    kyc_w = [0.15, 0.15, 0.62, 0.08]

    # Signups: mostly within month, some before month for more realistic active base
    rows = []
    for i in range(n_users):
        user_id = str(uuid.uuid4())
        # 75% sign up within the month, 25% earlier (up to 90 days)
        if rng.random() < 0.75:
            signup = month_start + timedelta(seconds=rng.randint(0, int((month_end - month_start).total_seconds()) - 1))
        else:
            signup = month_start - timedelta(days=rng.randint(1, 90), seconds=rng.randint(0, 86400))
        country = weighted_choice(rng, countries, country_w)
        device_os = weighted_choice(rng, oses, os_w)
        kyc_status = weighted_choice(rng, kyc, kyc_w)

        # Some accounts deleted (for event-after-delete anomaly checks)
        is_deleted = rng.random() < 0.04
        deleted_ts = None
        if is_deleted:
            # deleted after signup; often within month but can be later
            del_base = max(signup + timedelta(hours=1), month_start)
            deleted_ts = del_base + timedelta(days=rng.randint(0, 25), seconds=rng.randint(0, 86400))
            deleted_ts = clamp_dt(deleted_ts, month_start, month_end)

        rows.append(
            {
                "user_id": user_id,
                "signup_ts": ts_iso(signup),
                "country": country,
                "kyc_status": kyc_status,
                "device_os": device_os,
                "is_deleted": bool(is_deleted),
                "deleted_ts": ts_iso(deleted_ts) if deleted_ts else None,
            }
        )

    return pd.DataFrame(rows)


def pick_users_for_txn(rng: random.Random, user_ids: list[str]) -> tuple[str, str]:
    sender = rng.choice(user_ids)
    receiver = rng.choice(user_ids)
    # avoid self-transfer most of the time
    if receiver == sender and rng.random() < 0.9:
        while receiver == sender:
            receiver = rng.choice(user_ids)
    return sender, receiver


def generate_transactions(
    rng: random.Random,
    users: pd.DataFrame,
    month_start: datetime,
    month_end: datetime,
    cfg: Config,
) -> pd.DataFrame:
    user_ids = users["user_id"].tolist()
    signup_map = dict(zip(users["user_id"], pd.to_datetime(users["signup_ts"], utc=True)))

    currencies = ["GBP", "USD", "EUR"]
    currency_w = [0.55, 0.25, 0.20]
    statuses = ["succeeded", "failed", "reversed"]
    status_w = [0.92, 0.06, 0.02]
    channels = ["p2p_send", "p2p_request", "p2p_qr"]
    channel_w = [0.72, 0.18, 0.10]

    rows = []
    for _ in range(cfg.n_txns):
        sender, receiver = pick_users_for_txn(rng, user_ids)
        created = month_start + timedelta(seconds=rng.randint(0, int((month_end - month_start).total_seconds()) - 1))

        # Amounts: lognormal-ish with cents; clip to keep plausible
        amt = float(np.round(np.random.lognormal(mean=3.1, sigma=0.75), 2))  # ~22 median-ish
        amt = float(min(amt, 2500.00))

        currency = weighted_choice(rng, currencies, currency_w)
        status = weighted_choice(rng, statuses, status_w)
        channel = weighted_choice(rng, channels, channel_w)

        transaction_id = f"txn_{uuid.uuid4().hex[:16]}"

        # --- anomalies injected (different from brief examples) ---
        # 1) invalid amount (<= 0 or extreme outlier)
        if rng.random() < cfg.rate_txn_invalid_amount:
            amt = rng.choice([
                0.0,
                -1.0 * float(np.round(np.random.uniform(0.01, 50.0), 2)),
                250000.0,
            ])

        # 2) orphan receiver: receiver_user_id does not exist in users
        if rng.random() < cfg.rate_txn_orphan_receiver:
            receiver = str(uuid.uuid4())

        # 3) timestamp outside the reporting month window (but after sender signup)
        if rng.random() < cfg.rate_txn_outside_month:
            direction = -1 if rng.random() < 0.5 else 1
            created = created + timedelta(days=direction * rng.randint(1, 7))

        # make sure we still don't create the specific "before signup" anomaly
        signup_sender = signup_map[sender].to_pydatetime()
        if created < signup_sender:
            created = signup_sender + timedelta(minutes=rng.randint(1, 60))

        rows.append(
            {
                "transaction_id": transaction_id,
                "created_ts": ts_iso(created),
                "sender_user_id": sender,
                "receiver_user_id": receiver,
                "amount": amt,
                "currency": currency,
                "status": status,
                "channel": channel,
            }
        )

    txns = pd.DataFrame(rows)
    return txns


def generate_events(
    rng: random.Random,
    fake: Faker,
    users: pd.DataFrame,
    month_start: datetime,
    month_end: datetime,
    cfg: Config,
) -> pd.DataFrame:
    user_ids = users["user_id"].tolist()
    signup_map = dict(zip(users["user_id"], pd.to_datetime(users["signup_ts"], utc=True)))

    event_types = ["login", "page_view", "button_click", "p2p_initiate", "p2p_confirm"]
    event_w = [0.12, 0.55, 0.22, 0.06, 0.05]
    pages = ["/home", "/wallet", "/p2p", "/p2p/send", "/p2p/confirm", "/settings", "/activity"]
    buttons = ["send_cta", "request_cta", "qr_cta", "confirm_cta", "cancel_cta", None]
    btn_w = [0.18, 0.10, 0.05, 0.08, 0.04, 0.55]
    platforms = ["ios", "android", "web"]
    plat_w = [0.44, 0.46, 0.10]

    rows = []
    for _ in range(cfg.n_events):
        event_id = f"evt_{uuid.uuid4().hex[:18]}"
        event_ts = month_start + timedelta(seconds=rng.randint(0, int((month_end - month_start).total_seconds()) - 1))

        etype = weighted_choice(rng, event_types, event_w)
        platform = weighted_choice(rng, platforms, plat_w)
        page = rng.choice(pages) if etype in ["page_view", "p2p_initiate", "p2p_confirm"] else None
        button_id = weighted_choice(rng, buttons, btn_w) if etype in ["button_click", "p2p_initiate", "p2p_confirm"] else None
        session_id = f"s_{uuid.uuid4().hex[:12]}"

        # pick a user (or replace with unknown for anomalies)
        user_id = rng.choice(user_ids)

        # --- anomalies injected (different from brief examples) ---
        # unknown user id (referential integrity issue)
        if rng.random() < cfg.rate_event_unknown_user:
            user_id = str(uuid.uuid4())

        rows.append(
            {
                "event_id": event_id,
                "event_ts": ts_iso(event_ts),
                "user_id": user_id,
                "event_type": etype,
                "session_id": session_id,
                "page": page,
                "button_id": button_id,
                "platform": platform,
            }
        )

    events = pd.DataFrame(rows)
    return events


def build_charts(users: pd.DataFrame, txns: pd.DataFrame, events: pd.DataFrame, out_dir: str) -> dict[str, str]:
    ensure_dir(out_dir)

    txns_ok = txns.copy()
    txns_ok["created_ts"] = pd.to_datetime(txns_ok["created_ts"], utc=True, errors="coerce")
    txns_ok = txns_ok[(txns_ok["status"] == "succeeded") & (txns_ok["amount"].notna())]
    txns_ok = txns_ok[(txns_ok["amount"] > 0) & (txns_ok["amount"] <= 100000)]

    # Daily volume
    daily = (
        txns_ok.assign(day=txns_ok["created_ts"].dt.floor("D"))
        .groupby("day", as_index=False)["amount"]
        .sum()
        .rename(columns={"amount": "volume"})
    )
    fig = plt.figure(figsize=(8, 3))
    plt.plot(daily["day"], daily["volume"])
    plt.title("Daily P2P Volume (Succeeded, Clean-ish)")
    plt.xlabel("Day")
    plt.ylabel("Volume")
    plt.tight_layout()
    p1 = os.path.join(out_dir, "daily_volume.png")
    plt.savefig(p1, dpi=160)
    plt.close(fig)

    # DAU (events)
    ev = events.copy()
    ev["event_ts"] = pd.to_datetime(ev["event_ts"], utc=True, errors="coerce")
    ev = ev[ev["user_id"].notna()]
    dau = ev.assign(day=ev["event_ts"].dt.floor("D")).groupby("day")["user_id"].nunique().reset_index(name="dau")
    fig = plt.figure(figsize=(8, 3))
    plt.plot(dau["day"], dau["dau"])
    plt.title("Daily Active Users (Events, Raw)")
    plt.xlabel("Day")
    plt.ylabel("DAU")
    plt.tight_layout()
    p2 = os.path.join(out_dir, "dau.png")
    plt.savefig(p2, dpi=160)
    plt.close(fig)

    # Amount distribution
    fig = plt.figure(figsize=(6, 3))
    plt.hist(txns_ok["amount"].clip(upper=500), bins=40)
    plt.title("Transaction Amount Distribution (clipped at 500)")
    plt.xlabel("Amount")
    plt.ylabel("Count")
    plt.tight_layout()
    p3 = os.path.join(out_dir, "amount_hist.png")
    plt.savefig(p3, dpi=160)
    plt.close(fig)

    return {"daily_volume": p1, "dau": p2, "amount_hist": p3}


def build_pptx(
    users: pd.DataFrame,
    txns: pd.DataFrame,
    events: pd.DataFrame,
    charts: dict[str, str],
    out_path: str,
) -> None:
    prs = Presentation()

    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

    def add_chart_slide(title: str, img_path: str, caption: str | None = None) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        slide.shapes.title.text = title
        slide.shapes.add_picture(img_path, Inches(0.7), Inches(1.4), width=Inches(8.6))
        if caption:
            tx = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(8.6), Inches(0.4))
            tx.text_frame.text = caption

    # Headline metrics (clean-ish)
    tx = txns.copy()
    tx["created_ts"] = pd.to_datetime(tx["created_ts"], utc=True, errors="coerce")
    tx_cleanish = tx[(tx["status"] == "succeeded") & (tx["amount"].notna()) & (tx["amount"] > 0) & (tx["amount"] <= 100000)]
    total_volume = float(tx_cleanish["amount"].sum())
    active_users = int(events[events["user_id"].notna()]["user_id"].nunique())

    add_title_slide(
        "P2P Feature – 1 Month Performance (Synthetic)",
        "Audience: Data Team + Head of Product • Goal: validate analytics pipeline with realistic noisy data",
    )

    add_bullets_slide(
        "Executive Summary",
        [
            f"Total succeeded volume (clean-ish): {total_volume:,.0f}",
            f"Unique active users (raw events): {active_users:,}",
            "Volume shows steady activity with day-to-day variability (expected for consumer P2P)",
            "Data quality issues materially impact trusted metrics unless cleaned",
        ],
    )

    add_chart_slide("Daily Volume", charts["daily_volume"], "Succeeded txns with basic sanity filters; SQL provides the fully trusted version.")
    add_chart_slide("Daily Active Users (Raw)", charts["dau"], "Raw DAU can be inflated by bad user_ids; cleaning rules included in SQL.")
    add_chart_slide("Transaction Size Distribution", charts["amount_hist"], "Long tail is normal; invalid amounts are intentionally injected and must be filtered.")

    add_bullets_slide(
        "Data Health Report (What we injected & why)",
        [
            "Suspicious transaction amounts (<= 0 or extreme outliers) to test metric robustness",
            "Orphan transactions with receiver_user_id not present in the user base (foreign key drift)",
            "Events referencing non-existent user_id (client state corruption / bad identity binding)",
            "Transactions and events leaking slightly outside the reporting month window (pipeline windowing issues)",
        ],
    )

    add_bullets_slide(
        "Root Cause Fixes (Engineering)",
        [
            "Add range checks and business rules on amounts before events hit the warehouse",
            "Tighten referential constraints and reconciliation jobs for user_id / receiver_user_id fields",
            "Harden identity binding on the client and validate user_id at ingestion time",
            "Introduce clear reporting windows (watermarks) so off-window data is flagged and reviewed",
        ],
    )

    add_bullets_slide(
        "Strategic Recommendation",
        [
            "Scale cautiously: core engagement and volume look promising in cleaned metrics",
            "Prioritize data quality hardening before relying on metrics for product decisions",
            "Add monitoring: anomaly rate dashboards + alerts (duplicates, nulls, temporal violations)",
        ],
    )

    ensure_dir(os.path.dirname(out_path))
    prs.save(out_path)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--seed", type=int, default=42)
    parser.add_argument("--month", type=str, default="2026-01", help="YYYY-MM (data generated for this month window)")
    parser.add_argument("--out", type=str, default="data")
    args = parser.parse_args()

    rng = random.Random(args.seed)
    np.random.seed(args.seed)
    Faker.seed(args.seed)
    fake = Faker()

    month_start, month_end = month_window(args.month)
    cfg = Config()

    ensure_dir(args.out)
    ensure_dir("presentation")

    users = generate_users(rng, fake, month_start, month_end, cfg.n_users)
    txns = generate_transactions(rng, users, month_start, month_end, cfg)
    events = generate_events(rng, fake, users, month_start, month_end, cfg)

    users.to_csv(os.path.join(args.out, "users.csv"), index=False)
    txns.to_csv(os.path.join(args.out, "transactions.csv"), index=False)
    events.to_csv(os.path.join(args.out, "app_events.csv"), index=False)

    charts = build_charts(users, txns, events, out_dir=os.path.join("presentation", "charts"))
    build_pptx(users, txns, events, charts, out_path=os.path.join("presentation", "deck.pptx"))

    print("Wrote:")
    print(f"- {os.path.join(args.out, 'users.csv')}")
    print(f"- {os.path.join(args.out, 'transactions.csv')}")
    print(f"- {os.path.join(args.out, 'app_events.csv')}")
    print("- presentation/deck.pptx")


if __name__ == "__main__":
    main()

